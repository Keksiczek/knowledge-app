"""
text_extraction.py – converts uploaded files to plain text.

Supported formats:
  • TXT, MD  – read as-is (UTF-8)
  • PDF      – pypdf (text layer) + pytesseract OCR fallback
  • DOCX     – python-docx
  • PPTX     – python-pptx (slide text + notes)

Each extractor returns a single str.
"""
from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Callable

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────────────────────────────────────
# Individual extractors
# ──────────────────────────────────────────────────────────────────────────────

def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_pdf(path: Path) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            pages.append(text)
        full_text = "\n\n".join(pages)
        # If almost no text was extracted, try OCR
        if len(full_text.strip()) < 50:
            full_text = _ocr_pdf(path)
        return full_text
    except ImportError:
        logger.warning("pypdf not installed, falling back to OCR")
        return _ocr_pdf(path)


def _ocr_pdf(path: Path) -> str:
    """OCR via pdf2image + pytesseract (optional)."""
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(str(path), dpi=200)
        parts = [pytesseract.image_to_string(img) for img in images]
        return "\n\n".join(parts)
    except ImportError:
        logger.warning("pdf2image / pytesseract not installed; returning empty text for OCR PDF")
        return ""
    except Exception as exc:
        logger.error("OCR failed: %s", exc)
        return ""


def _extract_docx(path: Path) -> str:
    try:
        import docx  # python-docx
        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cell_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cell_texts:
                    paragraphs.append(" | ".join(cell_texts))
        return "\n\n".join(paragraphs)
    except ImportError:
        # fallback to docx2txt
        try:
            import docx2txt
            return docx2txt.process(str(path))
        except ImportError:
            logger.error("Neither python-docx nor docx2txt installed")
            return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # python-pptx

        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"## Slide {i}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            # Notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Notes] {notes}")
            slides.append("\n".join(parts))
        return "\n\n".join(slides)
    except ImportError:
        logger.error("python-pptx not installed")
        return ""


# ──────────────────────────────────────────────────────────────────────────────
# Dispatch table
# ──────────────────────────────────────────────────────────────────────────────

_EXTRACTORS: dict[str, Callable[[Path], str]] = {
    ".txt":  _extract_txt,
    ".md":   _extract_txt,
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
}

SUPPORTED_EXTENSIONS = set(_EXTRACTORS.keys())


def extract_text(file_path: Path) -> str:
    """Return the plain-text content of *file_path*.

    Raises ValueError for unsupported extensions.
    """
    ext = file_path.suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext}")
    try:
        text = extractor(file_path)
    except Exception as exc:
        logger.error("Extraction failed for %s: %s", file_path, exc)
        raise RuntimeError(f"Could not extract text from {file_path.name}: {exc}") from exc
    return text.strip()
